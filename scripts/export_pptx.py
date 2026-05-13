from pptx import Presentation
from pptx.util import Emu
from pathlib import Path

SLIDES_DIR = Path(__file__).parent.parent / "slides_export"
OUT_FILE = SLIDES_DIR / "firewatch_pitch.pptx"

# 16:9 — 1280x720 px @ 96dpi → inches → EMU
W = Emu(9144000)   # 10 inches
H = Emu(5143500)   # 5.625 inches

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # fully blank

slides = sorted(SLIDES_DIR.glob("slide_*.png"))
for img_path in slides:
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(str(img_path), 0, 0, width=W, height=H)
    print(f"  + {img_path.name}")

prs.save(str(OUT_FILE))
print(f"\n✓ Saved → {OUT_FILE}")
